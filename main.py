from tabnanny import check
from app import app
from flask import Flask, flash, request, redirect, render_template, jsonify
from werkzeug.utils import secure_filename
import cv2
import numpy as np
import io
from PIL import Image
import base64
from Helpers import *
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import language_tool_python
tool = language_tool_python.LanguageTool('en-US')

text = "Your the best but their are allso  good!"
matches = tool.check(text)
print(len(matches))
ALLOWED_EXTENSIONS = set(['ppt', 'pptx'])
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/')
def upload_form():
    return render_template('upload.html')

@app.route('/', methods=['POST'])
def upload_image():
    global i, n, images, text_runs
    i = 0
    n = 0
    images = []
    text_runs = []

    f = request.files.getlist("file[]")[0]
    prs = Presentation(f)
    iter_picture_shapes(prs)
    get_text(prs)
    print(images)
    #print(images[1][0])
    #base64img = getbase64_image(images[1][0])

    page = []
    for x, y in zip(text_runs, images):

        img1 = []
        text1 = []
        for k in x:

            check = tool.check(k)
            if len(check) > 0:
                text1.append(check)
        for k in y:
            print(k)
            img = cv2.imread(k)
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
            laplacian = cv2.Laplacian(gray, cv2.CV_64F)
            fm = laplacian.var()
            result = "Not Blurry"

            if fm < 250:
                result = "Blurry"
                print(result, fm)
                sharpness_value = "{:.0f}".format(fm)
                message = [result, sharpness_value]
                img1.append([message, getbase64_image(k)])
        if (len(img1) > 0 or len(text1) > 0):
            page.append([img1, text1])
        print(len(page))

    return render_template('upload.html', pages=page)


n = 0
i = 0
filename = 'test.pptx'
images = []
page = []


def write_image(shape):
    global n, i
    image = shape.image
    # ---get image "file" contents---
    image_bytes = image.blob
    # ---make up a name for the file, e.g. 'image.jpg'---
    image_filename = 'img/image{:03d}.{}'.format(n, image.ext)
    n += 1
    # print(image_filename)
    page.append(image_filename)
    with open(image_filename, 'wb') as f:
        f.write(image_bytes)


def visitor(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            visitor(s)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        write_image(shape)


def iter_picture_shapes(prs):
    global i, page
    for slide in prs.slides:
        page = []
        i += 1
        # print(i)
        for shape in slide.shapes:
            visitor(shape)
        images.append(page)


iter_picture_shapes(Presentation(filename))


def getbase64_image(image):
    img = img = cv2.imread(image)
    img = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
    file_object = io.BytesIO()
    img = Image.fromarray(Helpers.resize(img, width=500))
    img.save(file_object, 'PNG')
    return "data:image/png;base64,"+base64.b64encode(file_object.getvalue()).decode('ascii')

text_runs = []
prs = Presentation(filename)

def get_text(prs):
    for slide in prs.slides:
        text1 = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text1.append(run.text)
        text_runs.append(text1)

get_text(prs)

if __name__ == "__main__":
    app.run(debug=True)
print(images)
